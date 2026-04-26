"""Build the MediBill-Env pitch deck as a .pptx file (opens in Keynote).

Run:
    python3 scripts/build_deck.py

Output:
    docs/medibill_pitch.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent.parent / "docs" / "medibill_pitch.pptx"

# Brand-ish palette (kept simple — dark ink on white)
INK = RGBColor(0x10, 0x10, 0x10)
ACCENT = RGBColor(0x0E, 0x6B, 0xA8)        # blue
HIGHLIGHT = RGBColor(0x0A, 0x84, 0x3D)     # green for hero numbers
RULE = RGBColor(0xCC, 0xCC, 0xCC)


def add_title(slide, text):
    title = slide.shapes.title
    title.text = text
    for para in title.text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = INK


def add_bullets(slide, bullets):
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(22)
            run.font.color.rgb = INK


def add_notes(slide, notes):
    nf = slide.notes_slide.notes_text_frame
    nf.text = notes


def add_table(slide, left, top, width, height, headers, rows, *, hero_cells=None):
    rows_n = 1 + len(rows)
    cols_n = len(headers)
    tbl_shape = slide.shapes.add_table(rows_n, cols_n, left, top, width, height)
    tbl = tbl_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(16)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    # Body rows
    hero_cells = hero_cells or set()
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(14)
                    r.font.color.rgb = INK
                    if (i, j) in hero_cells:
                        r.font.bold = True
                        r.font.color.rgb = HIGHLIGHT


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_content = prs.slide_layouts[1]   # Title + content
    blank = prs.slide_layouts[5]           # Title only (we'll add tables)

    # ---- Slide 0: Cover -------------------------------------------------
    cover_layout = prs.slide_layouts[0]
    cover = prs.slides.add_slide(cover_layout)
    cover.shapes.title.text = "MediBill-Env"
    sub = cover.placeholders[1]
    sub.text = "Silent policy drift in Indian health-insurance claims\nMeta × Scaler OpenEnv Hackathon — Round 2"
    add_notes(cover, "Hi, I'm Anuj. MediBill-Env is an OpenEnv environment for testing whether an LLM agent can detect and recover from silent policy drift in medical claims billing.")

    # ---- Slide 1: The regulatory clock ---------------------------------
    s1 = prs.slides.add_slide(title_content)
    add_title(s1, "180 minutes to close the claim.")
    add_bullets(s1, [
        "IRDAI mandate (May 2024): 1 hour pre-auth, 3 hours discharge",
        "Miss the 3-hour clock → insurer eats the cost from shareholder funds",
        "FY24: ~₹26,000 cr health-claim disallowed",
        "~13% of pre-auths still miss the window",
    ])
    add_notes(s1, (
        "In India, IRDAI gives hospitals one hour for pre-authorization and three hours for final "
        "discharge on every cashless claim. Miss the three-hour clock, and the overrun comes out of "
        "the insurer's shareholder fund. Industry estimates put FY24 disallowed health-claim value "
        "around twenty-six thousand crore rupees. Roughly thirteen percent of pre-auths still miss "
        "the one-hour window. The bottleneck is a human coder racing a clock, and the policies keep "
        "changing on them."
    ))

    # ---- Slide 2: Problem is staleness ----------------------------------
    s2 = prs.slides.add_slide(title_content)
    add_title(s2, "Why agents fail here")
    add_bullets(s2, [
        "Rules engines handle static schema validation",
        "They do not handle staleness — yesterday's correct rule, today wrong",
        "Agents that imitate one month's trajectories fail quietly the next month",
        "We need an agent that knows to re-check before submitting",
    ])
    add_notes(s2, (
        "Most agent benchmarks check whether the agent can fill a form correctly. That is schema "
        "validation, and rules engines already do it. The real failure mode in this domain is "
        "staleness — the policy changed, the agent did not notice, the claim is wrong. An agent "
        "that learned by imitating last month's expert trajectories will reproduce last month's "
        "rules. We want an agent that knows to re-check before submitting."
    ))

    # ---- Slide 3: The environment ---------------------------------------
    s3 = prs.slides.add_slide(blank)
    add_title(s3, "MediBill-Env: 5 tools, 3 task tiers, 6-axis grader")
    # Tools table
    add_table(
        s3,
        left=Inches(0.6), top=Inches(1.4), width=Inches(6.5), height=Inches(3.5),
        headers=["Tool", "Purpose"],
        rows=[
            ["ehr_query", "Read patient record"],
            ["insurance_lookup", "Fetch active policy rules"],
            ["coding_engine", "Write a policy-sensitive field"],
            ["escalate_to_human", "Calibrated abstention"],
            ["submit_claim", "Lock claim for grading"],
        ],
    )
    # Tasks table
    add_table(
        s3,
        left=Inches(7.4), top=Inches(1.4), width=Inches(5.4), height=Inches(2.5),
        headers=["Task tier", "Drift?"],
        rows=[
            ["easy_cashless", "no"],
            ["medium_multi_payer", "no"],
            ["hard_drift", "yes — silent, mid-episode"],
        ],
    )
    # Footer
    txt = s3.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(12.0), Inches(1.0))
    tf = txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "6-axis deterministic grader · disjoint identity/policy partition asserted at import · 5 reward-hacking attacks neutralised"
    for r in p.runs:
        r.font.size = Pt(18)
        r.font.italic = True
        r.font.color.rgb = INK
    add_notes(s3, (
        "The agent has five tools: query the patient record, look up the insurer's active policy, "
        "write fields, escalate when uncertain, and submit. Three task tiers — easy, medium, and "
        "hard, where the policy mutates mid-episode. The grader has six axes with a disjoint field "
        "partition asserted at import time, so identity correctness and policy compliance never "
        "overlap."
    ))

    # ---- Slide 4: The hero mechanic -------------------------------------
    s4 = prs.slides.add_slide(title_content)
    add_title(s4, "Silent multi-field policy drift")
    add_bullets(s4, [
        "Active policy mutates 3–7 fields at a seed-randomized step",
        "No announcement — no observation flag, no metadata key, no event",
        "submit_claim is graded against the policy at submit time",
        "Only path to new rules: a fresh insurance_lookup after the drift step",
        "12 claim types × 3 tiers × randomized drift = ~12k+ unique trajectories",
        "Scripted baseline: 1.00 on easy, 0.7611 on drift — the 0.24 gap is the signal",
    ])
    add_notes(s4, (
        "On hard_drift tasks the active policy mutates mid-episode across three to seven fields — "
        "pre-auth thresholds, required signatures, narrative requirements, discharge attachment "
        "rules. Multi-field mutation, not a boolean. No announcement, no flag, no event. The only "
        "path to the new rules is a fresh insurance_lookup after the unknown drift step. "
        "Submissions are graded against the policy at submit time. Twelve claim types, three "
        "tiers, seed-randomized drift = over twelve thousand unique trajectories. Scripted "
        "baseline drops from one-zero on easy to zero-seven-six on drift. That zero-two-four gap "
        "is the trainable signal."
    ))

    # ---- Slide 5: HEADLINE — measurements -------------------------------
    s5 = prs.slides.add_slide(blank)
    add_title(s5, "Base 0.00 → SFT v2 0.9999 avg. Teacher engineering broke through GRPO saturation.")
    # Hero table — Base → SFT v2
    add_table(
        s5,
        left=Inches(0.5), top=Inches(1.3), width=Inches(6.3), height=Inches(2.6),
        headers=["task", "base Qwen", "SFT v2", "lift"],
        rows=[
            ["easy_cashless", "0.0000", "1.0000", "+1.000"],
            ["medium_multi_payer", "0.0000", "1.0000", "+1.000"],
            ["hard_drift", "0.0000", "0.9996 ± 0.0008", "+0.9996"],
            ["AVERAGE", "0.0000", "0.9999", "+0.9999"],
        ],
        hero_cells={(3, 2), (3, 3), (4, 1), (4, 2), (4, 3)},
    )
    # Iteration table
    add_table(
        s5,
        left=Inches(7.0), top=Inches(1.3), width=Inches(6.0), height=Inches(2.6),
        headers=["checkpoint", "hard_drift", "what changed"],
        rows=[
            ["Base Qwen 2.5 3B", "0.0000", "untrained"],
            ["SFT v1", "0.7573", "scripted teacher (parity)"],
            ["GRPO over SFT v1", "0.7575 (Δ±0.0002)", "rewards saturated — calibration"],
            ["SFT v2", "0.9996", "drift-aware teacher"],
        ],
        hero_cells={(4, 1), (4, 2)},
    )
    # Footer bullets
    foot = s5.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.5), Inches(2.5))
    tf = foot.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "• 5 exploit patterns explicitly neutralised — all five score ≤ no_op"
    p2 = tf.add_paragraph()
    p2.text = "• Pivot was teacher engineering, not RL — +0.2423 lift on hard_drift in 90 trajectories + 33 min retraining"
    p3 = tf.add_paragraph()
    p3.text = "• Verified via Codex reproducibility protocol: sha256 byte-match of adapter weights + fresh-subprocess re-eval × 2"
    for p in (p1, p2, p3):
        for r in p.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = INK
    add_notes(s5, (
        "Six bars on hard_drift, left to right: base Qwen at zero, random at eleven, no-op at "
        "eight, scripted at seventy-six, SFT v1 at seventy-six, our final SFT v2 at "
        "zero-point-nine-nine-nine-six. Untrained, the 3B model scores literal zero — zero parse "
        "failures across fifteen episodes — it can format JSON, it just has no policy reasoning. "
        "SFT v1 hit scripted-teacher parity. Then GRPO with five reward functions saturated — "
        "delta two ten-thousandths, gradient ten-to-minus-seven. Diagnosis: SFT extracts "
        "everything the rewards can grip on. So we engineered a stronger teacher — Scripted plus "
        "plus, which escalates ambiguous cells and does a fresh insurance lookup before each "
        "submit. Ninety new trajectories, thirty-three minutes of retraining. SFT v2: one-zero-zero "
        "on easy and medium, zero-point-nine-nine-nine-six on hard. Average lift base to SFT v2: "
        "zero-point-nine-nine-nine-nine."
    ))

    # ---- Slide 6: Scope + close -----------------------------------------
    s6 = prs.slides.add_slide(title_content)
    add_title(s6, "Environment-first submission under Theme 3.1")
    add_bullets(s6, [
        "Shipping today: env + grader + 5-attack exploit suite + scripted baseline + SFT v2 adapter (0.9999 avg)",
        "Two of six axes — abstention_quality and drift_bonus — are RL-only targets (spec v3 §7.6)",
        "Code enforces every claim: disjoint partition asserted at import, 5 exploit tests, prompt-version handshake",
        "Theme 3.1 — DataOps Copilot. Enterprise reasoning under shifting business rules.",
        "Repo: github.com/Algoace1403/METAHackthon2026",
        "HF Space (LIVE): huggingface.co/spaces/Anuj424614/medibill-env",
    ])
    add_notes(s6, (
        "We submit under Theme 3.1, DataOps Copilot. Shipping today: the environment, six-axis "
        "deterministic grader, silent drift mechanic, five-attack exploit suite, scripted "
        "baseline, and a trained SFT v2 adapter that hits zero-point-nine-nine-nine-nine average "
        "across all three difficulty tiers — table on slide five. Two axes — abstention and "
        "drift_bonus — are RL-only by design. Disjoint partition at import, five exploit tests, "
        "prompt-version handshake. Repo and live HF Space on screen. Thank you."
    ))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved deck to: {OUT}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
